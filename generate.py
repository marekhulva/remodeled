"""
ArchGram MVP — Layout Engine Test
Generates PPTX with N data center containers to test sizing/positioning.

Usage:
    python3 generate.py --sites 1
    python3 generate.py --sites 2
    python3 generate.py --sites 5
    python3 generate.py --sites 10

Output: output/test_N_sites.pptx
"""

import argparse
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from lxml import etree

# ============================================================
# DESIGN TOKENS (extracted from PPT Template)
# ============================================================

COLORS = {
    'bg':               RGBColor(0x00, 0x00, 0x00),  # #000000
    'purple_primary':   RGBColor(0x70, 0x30, 0xA0),  # #7030A0
    'purple_secondary': RGBColor(0x83, 0x48, 0x95),  # #834895
    'text_primary':     RGBColor(0xFF, 0xFF, 0xFF),  # #FFFFFF
    'text_muted':       RGBColor(0xE6, 0xE8, 0xF0),  # #E6E8F0
    'border_dark':      RGBColor(0x3C, 0x3F, 0x48),  # #3C3F48
    'border_medium':    RGBColor(0x5C, 0x5F, 0x6B),  # #5C5F6B
    'positive':         RGBColor(0x00, 0xB0, 0x50),  # #00B050
    'negative':         RGBColor(0xFF, 0x00, 0x00),  # #FF0000
}

FONT = 'Arial'

# Slide dimensions
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

# Layout margins
MARGIN_TOP = Inches(1.0)     # Space for title
MARGIN_BOTTOM = Inches(0.4)  # Space for status labels
MARGIN_LEFT = Inches(0.3)
MARGIN_RIGHT = Inches(0.3)
GAP_BETWEEN_SITES = Inches(0.3)  # Horizontal gap between site containers

# Usable area
USABLE_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
USABLE_HEIGHT = SLIDE_HEIGHT - MARGIN_TOP - MARGIN_BOTTOM

# ============================================================
# DETAIL LEVEL THRESHOLDS (in inches)
# ============================================================

FULL_MIN_WIDTH = 3.0       # Full detail: > 3" per site
REDUCED_MIN_WIDTH = 2.0    # Reduced detail: > 2" per site
COMPACT_MIN_WIDTH = 1.2    # Compact: > 1.2" per site
# Below compact → group or multi-slide


def determine_detail_level(site_count):
    """Determine the detail level based on how many sites need to fit."""
    usable_inches = USABLE_WIDTH / 914400 / 10  # Convert EMU to inches
    # Actually let's compute properly
    usable_in = 13.33 - 0.3 - 0.3  # slide width - margins
    total_gaps = GAP_BETWEEN_SITES / 914400 / 10 * (site_count - 1) if site_count > 1 else 0
    # Simpler calc
    usable_for_sites = 12.73  # 13.33 - 0.3 - 0.3
    gap_in = 0.3
    total_gap = gap_in * (site_count - 1) if site_count > 1 else 0
    width_per_site = (usable_for_sites - total_gap) / site_count

    if width_per_site >= FULL_MIN_WIDTH:
        return 'full', width_per_site
    elif width_per_site >= REDUCED_MIN_WIDTH:
        return 'reduced', width_per_site
    elif width_per_site >= COMPACT_MIN_WIDTH:
        return 'compact', width_per_site
    else:
        return 'grouped', width_per_site


def set_slide_bg_black(slide):
    """Set slide background to black."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['bg']


def make_transparent_table_with_border(slide, left, top, width, height, border_color=None):
    """
    Create a 1x1 table that acts as a transparent container with visible border.
    This matches the PPT template approach — tables used as bordered containers.
    """
    if border_color is None:
        border_color = COLORS['text_primary']

    # Add a 1x1 table
    table_shape = slide.shapes.add_table(1, 1, left, top, width, height)
    table = table_shape.table

    # Style the single cell
    cell = table.cell(0, 0)
    cell.text = ''

    # Make cell fill transparent (match background)
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # Remove any existing fill
    for existing in tcPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill'):
        tcPr.remove(existing)
    # Set no fill
    noFill = etree.SubElement(tcPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}noFill')

    # Set cell borders
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    borders = ['lnL', 'lnR', 'lnT', 'lnB']
    for border_name in borders:
        ln = etree.SubElement(tcPr, f'{{{ns}}}{border_name}')
        ln.set('w', str(int(Pt(1))))  # 1pt border width
        solidFill = etree.SubElement(ln, f'{{{ns}}}solidFill')
        srgbClr = etree.SubElement(solidFill, f'{{{ns}}}srgbClr')
        srgbClr.set('val', str(border_color))

    return table_shape


def add_site_label(slide, text, left, top, width, is_commvault=True):
    """Add site label above container with purple/grey underline."""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.28))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10.5)
    p.font.color.rgb = COLORS['text_primary']
    p.font.name = FONT
    p.font.bold = False

    # Add bottom border (purple underline effect) via shape line
    ln_color = COLORS['purple_primary'] if is_commvault else COLORS['border_dark']
    # Use lxml to add bottom border to the textbox shape
    sp = txBox._element
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    spPr = sp.find(f'{{{ns}}}spPr')
    if spPr is None:
        spPr = etree.SubElement(sp, f'{{{ns}}}spPr')
    ln = etree.SubElement(spPr, f'{{{ns}}}ln')
    ln.set('w', str(int(Pt(4.5))))  # 4.5pt border like template
    solidFill = etree.SubElement(ln, f'{{{ns}}}solidFill')
    srgbClr = etree.SubElement(solidFill, f'{{{ns}}}srgbClr')
    srgbClr.set('val', str(ln_color))

    return txBox


def add_slide_title(slide, text):
    """Add slide title at top."""
    txBox = slide.shapes.add_textbox(Inches(0.49), Inches(0.33), Inches(8.41), Inches(0.57))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['text_primary']
    p.font.name = FONT
    return txBox


# ============================================================
# FULL DETAIL RENDERER
# ============================================================

def render_full_detail(slide, site_name, left, top, width, height, site_data=None):
    """Render a site at full detail — container + all inner components."""
    label_height = Inches(0.28)
    label_gap = Inches(0.05)

    # Site label above container
    add_site_label(slide, site_name, left, top, width)

    # Main container (transparent + white border)
    container_top = top + label_height + label_gap
    container_height = height - label_height - label_gap
    make_transparent_table_with_border(slide, left, container_top, width, container_height)

    # Pull real data from site_data (or fall back to defaults)
    workloads = site_data.get('workloads', ['VMs']) if site_data else ['VMs']
    vm_count = site_data.get('vm_count', 100) if site_data else 100
    storage_tb = site_data.get('storage_tb', 10) if site_data else 10
    hsx_nodes = site_data.get('hsx_nodes', 3) if site_data else 3
    hsx_tb = site_data.get('hsx_tb', 150) if site_data else 150
    backup_target = site_data.get('backup_target', 'hsx') if site_data else 'hsx'
    is_commvault = (site_data.get('backup_software', 'commvault') if site_data else 'commvault') == 'commvault'

    # Inner padding
    pad = Inches(0.08)
    inner_left = left + pad
    inner_top = container_top + pad
    inner_width = width - pad * 2
    inner_y = inner_top

    # --- 1. Clients & Protected Storage sub-box ---
    bar_height = Inches(0.25)
    bar_color = COLORS['purple_primary'] if is_commvault else COLORS['border_dark']
    bar = slide.shapes.add_shape(1, inner_left, inner_y, inner_width, bar_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = bar_color
    bar.line.fill.background()
    tf = bar.text_frame
    p = tf.paragraphs[0]
    p.text = "Clients & Protected Storage"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_primary']
    p.font.name = FONT
    inner_y += bar_height

    # Workload chips — fit as many as possible
    chip_width = Inches(0.70)
    chip_height = Inches(0.40)
    chip_gap = Inches(0.05)
    chip_y = inner_y + Inches(0.08)

    for i, label in enumerate(workloads):
        cx = inner_left + Inches(0.05) + (chip_width + chip_gap) * i
        if cx + chip_width > inner_left + inner_width:
            break
        chip = slide.shapes.add_shape(1, cx, chip_y, chip_width, chip_height)
        chip.fill.solid()
        chip.fill.fore_color.rgb = COLORS['bg']
        chip.line.color.rgb = COLORS['border_dark']
        chip.line.width = Pt(0.75)
        tf = chip.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(8)
        p.font.color.rgb = COLORS['text_primary']
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER

    inner_y += chip_height + Inches(0.15)

    # Summary line
    summary = slide.shapes.add_textbox(inner_left, inner_y, inner_width, Inches(0.25))
    tf = summary.text_frame
    p = tf.paragraphs[0]
    p.text = f"{vm_count} VMs | {storage_tb}TB"
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS['text_primary']
    p.font.name = FONT
    inner_y += Inches(0.35)

    # --- 2. Backup Software Stack ---
    stack_height = Inches(0.88)
    stack_bg = slide.shapes.add_shape(1, inner_left, inner_y, inner_width, stack_height)
    stack_bg.fill.solid()
    stack_bg.fill.fore_color.rgb = RGBColor(0x1A, 0x0A, 0x2E)  # Very dark purple
    stack_bg.line.color.rgb = COLORS['purple_secondary']
    stack_bg.line.width = Pt(1)

    if is_commvault:
        # Commvault logo image
        cv_logo_path = os.path.join(os.path.dirname(__file__), 'assets/extracted/slide6_Picture_5_0accdd8f.png')
        logo_h = Inches(0.35)
        logo_w = Inches(1.80)
        logo_left = inner_left + Inches(0.05)
        slide.shapes.add_picture(cv_logo_path, logo_left, inner_y + Inches(0.05), logo_w, logo_h)

        # "Command Center" label below logo
        cc_label = slide.shapes.add_textbox(
            inner_left + Inches(0.05), inner_y + logo_h + Inches(0.08),
            inner_width - Inches(0.5), Inches(0.20)
        )
        tf = cc_label.text_frame
        p = tf.paragraphs[0]
        p.text = "Command Center"
        p.font.size = Pt(8)
        p.font.color.rgb = COLORS['text_primary']
        p.font.name = FONT

        # CS badge (purple oval) — right side
        badge_size = Inches(0.35)
        badge = slide.shapes.add_shape(
            9,  # OVAL
            inner_left + inner_width - badge_size - Inches(0.08),
            inner_y + (stack_height - badge_size) / 2,
            badge_size, badge_size
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = COLORS['purple_primary']
        badge.line.color.rgb = COLORS['purple_secondary']
        badge.line.width = Pt(1)
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = "CS"
        p.font.size = Pt(7)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_primary']
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER

    inner_y += stack_height + Inches(0.1)

    # --- 3. Backup Target sub-box ---
    target_bar = slide.shapes.add_shape(1, inner_left, inner_y, inner_width, bar_height)
    target_bar.fill.solid()
    target_bar.fill.fore_color.rgb = bar_color
    target_bar.line.fill.background()
    tf = target_bar.text_frame
    p = tf.paragraphs[0]
    p.text = "Protected Data Layer"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_primary']
    p.font.name = FONT
    inner_y += bar_height

    if backup_target == 'pure':
        # Pure Storage logo
        pure_path = os.path.join(os.path.dirname(__file__), 'assets/extracted/slide5_Picture_4_788cf27a.png')
        logo_w = Inches(1.39)
        logo_h = Inches(0.53)
        logo_left = inner_left + (inner_width - logo_w) / 2
        slide.shapes.add_picture(pure_path, logo_left, inner_y + Inches(0.1), logo_w, logo_h)
        inner_y += logo_h + Inches(0.2)

    else:
        # HSX table — 2 rows per node
        total_rows = hsx_nodes * 2
        row_height = Inches(0.12)
        hsx_height = row_height * total_rows

        for r in range(total_rows):
            row_y = inner_y + row_height * r
            node_num = r // 2 + 1
            is_label_row = r % 2 == 0

            slot = slide.shapes.add_shape(1, inner_left + Inches(0.05), row_y, Inches(0.25), row_height)
            slot.fill.solid()
            slot.fill.fore_color.rgb = COLORS['purple_primary']
            slot.line.color.rgb = COLORS['border_dark']
            slot.line.width = Pt(0.5)
            tf = slot.text_frame
            p = tf.paragraphs[0]
            p.text = str(r + 1)
            p.font.size = Pt(6)
            p.font.color.rgb = COLORS['text_primary']
            p.font.name = FONT
            p.alignment = PP_ALIGN.CENTER

            node = slide.shapes.add_shape(1, inner_left + Inches(0.30), row_y, Inches(0.80), row_height)
            node.fill.solid()
            node.fill.fore_color.rgb = COLORS['purple_primary']
            node.line.color.rgb = COLORS['border_dark']
            node.line.width = Pt(0.5)
            tf = node.text_frame
            p = tf.paragraphs[0]
            p.text = f"HSX - {node_num:02d}" if is_label_row else ""
            p.font.size = Pt(6)
            p.font.color.rgb = COLORS['text_primary']
            p.font.name = FONT

        inner_y += hsx_height + Inches(0.05)

        hsx_label = slide.shapes.add_textbox(inner_left, inner_y, inner_width, Inches(0.25))
        tf = hsx_label.text_frame
        p = tf.paragraphs[0]
        p.text = f"{hsx_nodes}-Node HSX Appliance | {hsx_tb}TB Usable"
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_muted']
        p.font.name = FONT
        inner_y += Inches(0.25)

    # Protection status
    status = slide.shapes.add_textbox(inner_left, inner_y, inner_width, Inches(0.40))
    tf = status.text_frame
    p = tf.paragraphs[0]
    p.text = "Immutable | Deduped | Encrypted"
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS['text_primary']
    p.font.name = FONT


# ============================================================
# REDUCED DETAIL RENDERER
# ============================================================

def render_reduced_detail(slide, site_name, left, top, width, height, site_data=None):
    """Render a site at reduced detail — container + summary + backup name."""
    label_height = Inches(0.28)
    label_gap = Inches(0.05)

    add_site_label(slide, site_name, left, top, width)

    container_top = top + label_height + label_gap
    container_height = height - label_height - label_gap
    make_transparent_table_with_border(slide, left, container_top, width, container_height)

    vm_count = site_data.get('vm_count', 100) if site_data else 100
    storage_tb = site_data.get('storage_tb', 10) if site_data else 10
    hsx_nodes = site_data.get('hsx_nodes', 3) if site_data else 3
    hsx_tb = site_data.get('hsx_tb', 150) if site_data else 150

    pad = Inches(0.1)
    inner_left = left + pad
    inner_width = width - pad * 2
    inner_y = container_top + pad

    # Workload summary (no individual chips)
    summary = slide.shapes.add_textbox(inner_left, inner_y, inner_width, Inches(0.25))
    tf = summary.text_frame
    p = tf.paragraphs[0]
    p.text = f"{vm_count} VMs | {storage_tb}TB"
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS['text_primary']
    p.font.name = FONT
    inner_y += Inches(0.35)

    # Backup product name
    bp = slide.shapes.add_textbox(inner_left, inner_y, inner_width, Inches(0.20))
    tf = bp.text_frame
    p = tf.paragraphs[0]
    p.text = "Commvault"
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS['purple_secondary']
    p.font.name = FONT
    p.font.bold = True
    inner_y += Inches(0.30)

    # CS badge small
    badge_size = Inches(0.28)
    badge = slide.shapes.add_shape(9, inner_left, inner_y, badge_size, badge_size)
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLORS['purple_primary']
    badge.line.color.rgb = COLORS['purple_secondary']
    badge.line.width = Pt(1)
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = "CS"
    p.font.size = Pt(5)
    p.font.color.rgb = COLORS['text_primary']
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    inner_y += Inches(0.40)

    # HSX summary (no table)
    hsx = slide.shapes.add_textbox(inner_left, inner_y, inner_width, Inches(0.20))
    tf = hsx.text_frame
    p = tf.paragraphs[0]
    p.text = f"{hsx_nodes}-Node HSX | {hsx_tb}TB"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_muted']
    p.font.name = FONT
    inner_y += Inches(0.30)

    # Protection status
    status = slide.shapes.add_textbox(inner_left, inner_y, inner_width, Inches(0.30))
    tf = status.text_frame
    p = tf.paragraphs[0]
    p.text = "Immutable | Deduped | Encrypted"
    p.font.size = Pt(8)
    p.font.color.rgb = COLORS['text_primary']
    p.font.name = FONT


# ============================================================
# COMPACT RENDERER
# ============================================================

def render_compact(slide, site_name, left, top, width, height, site_data=None):
    """Render a site at compact detail — small box with key stats."""
    label_height = Inches(0.22)
    label_gap = Inches(0.03)

    # Smaller label
    txBox = slide.shapes.add_textbox(left, top, width, label_height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = site_name
    p.font.size = Pt(8)
    p.font.color.rgb = COLORS['text_primary']
    p.font.name = FONT
    p.font.bold = True

    container_top = top + label_height + label_gap
    container_height = height - label_height - label_gap
    make_transparent_table_with_border(slide, left, container_top, width, container_height)

    vm_count = site_data.get('vm_count', 100) if site_data else 100
    storage_tb = site_data.get('storage_tb', 10) if site_data else 10

    pad = Inches(0.08)
    inner_left = left + pad
    inner_width = width - pad * 2
    inner_y = container_top + pad

    # VM count
    vm = slide.shapes.add_textbox(inner_left, inner_y, inner_width, Inches(0.18))
    tf = vm.text_frame
    p = tf.paragraphs[0]
    p.text = f"{vm_count} VMs | {storage_tb}TB"
    p.font.size = Pt(7)
    p.font.color.rgb = COLORS['text_primary']
    p.font.name = FONT
    inner_y += Inches(0.22)

    # Vendor badge
    badge_size = Inches(0.25)
    badge = slide.shapes.add_shape(9, inner_left, inner_y, badge_size, badge_size)
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLORS['purple_primary']
    badge.line.color.rgb = COLORS['purple_secondary']
    badge.line.width = Pt(1)
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = "CV"
    p.font.size = Pt(5)
    p.font.color.rgb = COLORS['text_primary']
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER


# ============================================================
# GROUPED RENDERER
# ============================================================

def render_grouped(slide, sites, left, top, width, height):
    """Render multiple sites as a single grouped box."""
    make_transparent_table_with_border(slide, left, top, width, height)

    # Group label
    label = slide.shapes.add_textbox(
        left + Inches(0.1), top + Inches(0.1),
        width - Inches(0.2), Inches(0.30)
    )
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = f"{len(sites)} Data Centers"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['text_primary']
    p.font.name = FONT
    p.font.bold = True

    # List site names
    y = top + Inches(0.5)
    for site_name in sites:
        item = slide.shapes.add_textbox(left + Inches(0.15), y, width - Inches(0.3), Inches(0.20))
        tf = item.text_frame
        p = tf.paragraphs[0]
        p.text = f"• {site_name}"
        p.font.size = Pt(8)
        p.font.color.rgb = COLORS['text_muted']
        p.font.name = FONT
        y += Inches(0.22)


# ============================================================
# MAIN LAYOUT ENGINE
# ============================================================

def generate_diagram(site_count=None, scenario=None):
    """Generate a PPTX. Pass scenario dict (from JSON) or just a site_count for dummy data."""

    # Build site list
    if scenario:
        sites = scenario['sites']
        title = scenario.get('title', f"Future State — {len(sites)} Sites")
        site_count = len(sites)
    else:
        sites = []
        for i in range(site_count):
            name = "Primary DC" if i == 0 else ("DR Site" if i == 1 else f"Site {i + 1}")
            sites.append({'name': name, 'vm_count': 100, 'storage_tb': 10,
                          'workloads': ['VMs'], 'backup_software': 'commvault',
                          'hsx_nodes': 3, 'hsx_tb': 150})
        title = f"Future State — {site_count} Site{'s' if site_count > 1 else ''}"

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    detail_level, width_per_site = determine_detail_level(site_count)

    print(f"Sites: {site_count}")
    print(f"Detail level: {detail_level}")
    print(f"Width per site: {width_per_site:.2f} inches")

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg_black(slide)
    add_slide_title(slide, title)

    if detail_level == 'grouped':
        print("→ Using grouped mode (too many for one slide)")
        site_names = [s['name'] for s in sites]
        group_width = Inches(6)
        group_height = Inches(min(5.0, 1.0 + site_count * 0.25))
        group_left = (SLIDE_WIDTH - group_width) // 2
        group_top = MARGIN_TOP + Inches(0.5)
        render_grouped(slide, site_names, group_left, group_top, group_width, group_height)

        for i in range(0, site_count, 2):
            detail_slide = prs.slides.add_slide(slide_layout)
            set_slide_bg_black(detail_slide)
            batch = sites[i:i+2]
            add_slide_title(detail_slide, f"Detail — {', '.join(s['name'] for s in batch)}")
            w = Inches(5.5)
            for j, site in enumerate(batch):
                x = MARGIN_LEFT + (w + GAP_BETWEEN_SITES) * j
                render_full_detail(detail_slide, site['name'], x, MARGIN_TOP, w, USABLE_HEIGHT, site)

    else:
        site_width = Inches(width_per_site)
        renderer = {
            'full': render_full_detail,
            'reduced': render_reduced_detail,
            'compact': render_compact,
        }[detail_level]

        for i, site in enumerate(sites):
            x = MARGIN_LEFT + (site_width + GAP_BETWEEN_SITES) * i
            renderer(slide, site['name'], x, MARGIN_TOP, site_width, USABLE_HEIGHT, site)

    os.makedirs('output', exist_ok=True)
    filename = f'output/test_{site_count}_sites.pptx'
    prs.save(filename)
    print(f"Saved: {filename}")
    return filename


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='ArchGram Layout Engine Test')
    parser.add_argument('--sites', type=int, default=None, help='Number of dummy sites')
    parser.add_argument('--input', type=str, default=None, help='Path to scenario JSON file')
    args = parser.parse_args()

    if args.input:
        with open(args.input) as f:
            scenario = json.load(f)
        generate_diagram(scenario=scenario)
    elif args.sites:
        generate_diagram(site_count=args.sites)
    else:
        print("Usage: python3 generate.py --input scenario.json  OR  --sites N")

