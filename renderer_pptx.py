"""
PPTX Renderer — takes positioned shapes JSON from layout_engine → generates PPTX
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.dml.color import RGBColor
from lxml import etree

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
SLIDE_W = 13.33
SLIDE_H = 7.5

ALIGN_MAP = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
VALIGN_MAP = {'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM}


def rgb(hex_color):
    if not hex_color:
        return None
    h = hex_color.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def apply_gradient(shape, from_hex, to_hex):
    """Replace a shape's fill with a vertical linear gradient (top → bottom)."""
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    spPr = shape.fill._xPr  # ShapeProperties XML
    # Remove any existing fill children
    for tag in ('solidFill', 'gradFill', 'blipFill', 'pattFill', 'noFill'):
        for el in spPr.findall(f'{{{ns}}}{tag}'):
            spPr.remove(el)
    gradFill = etree.SubElement(spPr, f'{{{ns}}}gradFill',
                                attrib={'rotWithShape': '1', 'flip': 'none'})
    gsLst = etree.SubElement(gradFill, f'{{{ns}}}gsLst')
    for pos, color in ((0, from_hex), (100000, to_hex)):
        gs = etree.SubElement(gsLst, f'{{{ns}}}gs', attrib={'pos': str(pos)})
        srgb = etree.SubElement(gs, f'{{{ns}}}srgbClr',
                                attrib={'val': color.lstrip('#')})
    # 90° = top to bottom (angle in 60000ths of a degree)
    etree.SubElement(gradFill, f'{{{ns}}}lin',
                     attrib={'ang': '5400000', 'scaled': '0'})
    # Tile rect needs to be present to avoid warnings
    etree.SubElement(gradFill, f'{{{ns}}}tileRect')


def transparent_bordered_rect(slide, x, y, w, h, stroke):
    """Use 1x1 table trick for transparent container with visible border."""
    tbl = slide.shapes.add_table(1, 1, x, y, w, h).table
    cell = tbl.cell(0, 0)
    cell.text = ''
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for el in tcPr.findall(f'{{{ns}}}solidFill'):
        tcPr.remove(el)
    etree.SubElement(tcPr, f'{{{ns}}}noFill')
    for border in ['lnL', 'lnR', 'lnT', 'lnB']:
        ln = etree.SubElement(tcPr, f'{{{ns}}}{border}')
        ln.set('w', str(int(Pt(1))))
        sf = etree.SubElement(ln, f'{{{ns}}}solidFill')
        clr = etree.SubElement(sf, f'{{{ns}}}srgbClr')
        clr.set('val', stroke.lstrip('#'))


def draw(slide, s):
    t = s['type']

    if t == 'line':
        x1, y1 = Inches(s['x1']), Inches(s['y1'])
        x2, y2 = Inches(s['x2']), Inches(s['y2'])
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        conn.line.color.rgb = rgb(s.get('stroke', '#5C5F6B'))
        conn.line.width = Pt(s.get('sw', 1))
        dash = s.get('dash', 'dash')
        if dash == 'dash':
            conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        elif dash == 'dot':
            conn.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
        # Arrowheads — set tailEnd / headEnd on the connector's <a:ln>
        arrow = s.get('arrow')
        if arrow:
            ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            spPr = conn.line._get_or_add_ln()
            for end_tag, flag in (('tailEnd', arrow in ('end', 'both')),
                                  ('headEnd', arrow in ('start', 'both'))):
                if not flag:
                    continue
                # Remove any existing same-tag end markers
                for el in spPr.findall(f'{{{ns}}}{end_tag}'):
                    spPr.remove(el)
                end_el = etree.SubElement(spPr, f'{{{ns}}}{end_tag}')
                end_el.set('type', 'triangle')
                end_el.set('w', 'med')
                end_el.set('len', 'med')
        return

    x, y, w, h = Inches(s['x']), Inches(s['y']), Inches(s['w']), Inches(s['h'])

    if t == 'rect':
        radius = s.get('radius', 0) or 0
        gradient = s.get('gradient')
        if s.get('fill') is None and s.get('stroke') and radius == 0 and not gradient:
            transparent_bordered_rect(slide, x, y, w, h, s['stroke'])
        else:
            # MSO_SHAPE: 1 = RECTANGLE, 5 = ROUNDED_RECTANGLE
            shape_id = 5 if radius > 0 else 1
            shape = slide.shapes.add_shape(shape_id, x, y, w, h)
            if gradient:
                apply_gradient(shape, gradient[0], gradient[1])
            elif s.get('fill'):
                shape.fill.solid()
                shape.fill.fore_color.rgb = rgb(s['fill'])
            else:
                shape.fill.background()
            if s.get('stroke'):
                shape.line.color.rgb = rgb(s['stroke'])
                shape.line.width = Pt(s.get('sw', 1))
            else:
                shape.line.fill.background()

    elif t == 'text':
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = VALIGN_MAP.get(s.get('valign', 'top'), MSO_ANCHOR.TOP)
        # Remove default text frame margins so vertical centering is tight
        tf.margin_top = Pt(0)
        tf.margin_bottom = Pt(0)
        tf.margin_left = Pt(2)
        tf.margin_right = Pt(2)
        p = tf.paragraphs[0]
        p.text = s.get('text', '')
        p.font.size = Pt(s.get('fs', 10))
        p.font.color.rgb = rgb(s.get('color', '#FFFFFF'))
        p.font.name = 'Arial'
        p.font.bold = s.get('bold', False)
        p.alignment = ALIGN_MAP.get(s.get('align', 'left'), PP_ALIGN.LEFT)

    elif t == 'oval':
        shape = slide.shapes.add_shape(9, x, y, w, h)
        if s.get('fill'):
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(s['fill'])
        if s.get('stroke'):
            shape.line.color.rgb = rgb(s['stroke'])
            shape.line.width = Pt(s.get('sw', 1))
        if s.get('text'):
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = s['text']
            p.font.size = Pt(s.get('fs', 7))
            p.font.color.rgb = rgb(s.get('text_color', '#FFFFFF'))
            p.font.name = 'Arial'
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

    elif t == 'image':
        img_path = os.path.join(BASE_DIR, 'assets', s['src'])
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, x, y, w, h)


SLIDE_PAD = 0.2  # extra inch padding inside the slide before content


def _scale_shapes(shapes, scale, ox, oy):
    """Return new shape dicts with all coordinates scaled by `scale` and
    translated by (ox, oy). Lines have x1/y1/x2/y2; everything else has
    x/y/w/h. Used to fit a wide canvas onto a 13.33×7.5 slide."""
    out = []
    for s in shapes:
        n = dict(s)
        if s['type'] == 'line':
            n['x1'] = s['x1'] * scale + ox
            n['y1'] = s['y1'] * scale + oy
            n['x2'] = s['x2'] * scale + ox
            n['y2'] = s['y2'] * scale + oy
        else:
            n['x'] = s.get('x', 0) * scale + ox
            n['y'] = s.get('y', 0) * scale + oy
            n['w'] = s.get('w', 0) * scale
            n['h'] = s.get('h', 0) * scale
            if 'fs' in s:
                n['fs'] = max(4, s['fs'] * scale)  # shrink font with content
            if 'radius' in s:
                n['radius'] = s.get('radius', 0) * scale
            if 'sw' in s:
                n['sw'] = max(0.25, s['sw'] * scale)
        out.append(n)
    return out


def render_pptx(layout_data, output_path):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, layout_data['background'])

    # Fit-to-slide: scale uniformly so the diagram's content_w × content_h
    # fits inside (SLIDE_W - 2*SLIDE_PAD) × (SLIDE_H - 2*SLIDE_PAD), then
    # center the result. If content already fits, scale = 1.0 (untouched).
    content_w = layout_data.get('content_w', SLIDE_W)
    content_h = layout_data.get('content_h', SLIDE_H)
    avail_w = SLIDE_W - 2 * SLIDE_PAD
    avail_h = SLIDE_H - 2 * SLIDE_PAD
    scale = min(avail_w / content_w, avail_h / content_h, 1.0)
    final_w = content_w * scale
    final_h = content_h * scale
    ox = (SLIDE_W - final_w) / 2
    oy = (SLIDE_H - final_h) / 2

    shapes = (layout_data['shapes'] if scale == 1.0
              else _scale_shapes(layout_data['shapes'], scale, ox, oy))
    for s in shapes:
        draw(slide, s)

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    return output_path
